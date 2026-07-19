"""文档解析 + 智能分块"""

import os
import uuid
import re
from pathlib import Path
from app.core.config import get_settings

settings = get_settings()


def parse_document(file_path: str) -> str:
    """
    多格式文档解析 → 纯文本
    支持: PDF / PPT / PPTX / DOC / DOCX / MD / TXT
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in (".ppt", ".pptx"):
        return _parse_ppt(file_path)
    elif ext == ".docx":
        return _parse_docx(file_path)
    elif ext == ".doc":
        # 旧格式 .doc 不支持，提示用户另存为 .docx
        raise ValueError("不支持旧版 .doc 格式，请在 Word 中另存为 .docx 格式后再上传")
    elif ext in (".md", ".txt"):
        return Path(file_path).read_text(encoding="utf-8")
    else:
        raise ValueError(f"不支持的文件格式: {ext}")


def _parse_pdf(file_path: str) -> str:
    """PDF 解析，保留页码信息"""
    import fitz  # PyMuPDF
    doc = fitz.open(file_path)
    texts = []
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text()
        if text.strip():
            texts.append(f"[PAGE:{page_num}]\n{text}")
    return "\n\n".join(texts)


def _parse_ppt(file_path: str) -> str:
    """PPT 解析"""
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        slide_texts.append(para.text)
        if slide_texts:
            texts.append(f"[SLIDE:{slide_num}]\n" + "\n".join(slide_texts))
    return "\n\n".join(texts)


def _parse_docx(file_path: str) -> str:
    """Word 文档解析"""
    from docx import Document
    doc = Document(file_path)
    texts = []
    for para in doc.paragraphs:
        if para.text.strip():
            # 检测标题样式
            if para.style.name.startswith("Heading"):
                level = para.style.name.split()[-1]
                texts.append(f"{'#' * int(level)} {para.text}")
            else:
                texts.append(para.text)
    return "\n\n".join(texts)


def smart_chunk(text: str, chunk_size: int = None, overlap: int = None) -> list[dict]:
    """
    智能分块 —— 按语义边界切分，而非简单按字数切断

    策略:
    1. 先按页面/幻灯片分隔符分段
    2. 每段内按段落/标题边界切分
    3. 相邻块保留 overlap 重叠，避免切断语义
    """
    chunk_size = chunk_size or settings.chunk_size
    overlap = overlap or settings.chunk_overlap

    # 按页面分隔符拆分，保留页面信息
    page_pattern = re.compile(r'\[PAGE:(\d+)\]\n')
    slide_pattern = re.compile(r'\[SLIDE:(\d+)\]\n')

    # 统一用分段符标记；标记前后都用双换行隔离，使其独占一段——
    # 否则标记与正文粘在同一段，re.match 命中标记后 continue 会把正文一起丢弃
    text = page_pattern.sub(lambda m: f"\n\n{{PAGE:{m.group(1)}}}\n\n", text)
    text = slide_pattern.sub(lambda m: f"\n\n{{SLIDE:{m.group(1)}}}\n\n", text)

    paragraphs = text.split("\n\n")
    chunks = []
    current_chunk = ""
    current_page = 0
    chunk_page = 0  # 当前 chunk 的起始页码（引用标注"第 X 页"用）
    chunk_index = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        # 检测页面标记
        page_match = re.match(r'\{PAGE:(\d+)\}', para)
        slide_match = re.match(r'\{SLIDE:(\d+)\}', para)
        if page_match:
            current_page = int(page_match.group(1))
            continue
        if slide_match:
            current_page = int(slide_match.group(1))  # 幻灯片的"页码"
            continue

        # 如果加上当前段落会超过 chunk_size，则保存当前 chunk
        if len(current_chunk) + len(para) > chunk_size and current_chunk:
            chunks.append({
                "id": str(uuid.uuid4()),
                "content": current_chunk.strip(),
                "chunk_index": chunk_index,
                "page_number": chunk_page if chunk_page > 0 else None,
            })
            chunk_index += 1
            # 保留 overlap 部分
            overlap_text = current_chunk[-overlap:] if len(current_chunk) > overlap else current_chunk
            current_chunk = overlap_text + "\n\n" + para
            chunk_page = current_page  # 新 chunk 以当前段落页码为起始
        else:
            if current_chunk:
                current_chunk += "\n\n" + para
            else:
                current_chunk = para
                chunk_page = current_page  # 记录 chunk 起始页

    # 保存最后一个 chunk
    if current_chunk.strip():
        chunks.append({
            "id": str(uuid.uuid4()),
            "content": current_chunk.strip(),
            "chunk_index": chunk_index,
            "page_number": chunk_page if chunk_page > 0 else None,
        })

    return chunks
